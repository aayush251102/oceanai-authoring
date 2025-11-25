from fastapi import APIRouter, Depends, HTTPException
from fastapi.responses import FileResponse
from sqlalchemy.orm import Session
import json
import os
from pydantic import BaseModel

from ..database import SessionLocal
from ..models import Project
from ..auth import get_current_user
from ..llm import generate_section, refine_content

from docx import Document
from pptx import Presentation

router = APIRouter(prefix="/content")


# -------------------------
# DB SESSION
# -------------------------
def get_db():
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()


# -------------------------
# GENERATE FULL DOCUMENT
# -------------------------
@router.post("/{project_id}/generate-content")
def generate_content(project_id: int, token: str, db: Session = Depends(get_db)):
    user = get_current_user(token, db)

    project = db.query(Project).filter(Project.id == project_id).first()
    if not project or project.user_id != user.id:
        raise HTTPException(status_code=404, detail="Project not found")

    outline = json.loads(project.outline)
    topic = project.topic

    generated = {}
    for section in outline:
        generated[section] = generate_section(topic, section)

    project.content = json.dumps(generated)
    db.commit()

    return {"message": "Content generated", "content": generated}


# -------------------------
# REFINEMENT MODEL
# -------------------------
class RefineModel(BaseModel):
    section: str
    instruction: str


# ✅ AI REFINE (AUTO UPDATE)
@router.post("/{project_id}/refine-section")
def refine_section(
    project_id: int,
    body: RefineModel,
    token: str,
    db: Session = Depends(get_db),
):
    user = get_current_user(token, db)

    project = db.query(Project).filter(Project.id == project_id).first()
    if not project or project.user_id != user.id:
        raise HTTPException(status_code=404, detail="Project not found")

    content = json.loads(project.content)
    history = json.loads(project.history)

    if body.section not in content:
        raise HTTPException(status_code=400, detail="Section not found")

    old_text = content[body.section]
    new_text = refine_content(old_text, body.instruction)

    content[body.section] = new_text
    project.content = json.dumps(content)

    if body.section not in history:
        history[body.section] = []

    history[body.section].append({
        "old_text": old_text,
        "new_text": new_text,
        "instruction": body.instruction
    })

    project.history = json.dumps(history)
    db.commit()

    return {"message": "Section refined", "new_text": new_text}

# -------------------------
# GET HISTORY
# -------------------------
@router.get("/{project_id}/history")
def get_history(project_id: int, token: str, db: Session = Depends(get_db)):
    user = get_current_user(token, db)

    project = db.query(Project).filter(Project.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    return json.loads(project.history)


# -------------------------
# FEEDBACK MODEL
# -------------------------
class FeedbackModel(BaseModel):
    section: str
    feedback: str   # like / dislike


@router.post("/{project_id}/feedback")
def add_feedback(project_id: int, data: FeedbackModel, token: str, db: Session = Depends(get_db)):
    user = get_current_user(token, db)

    project = db.query(Project).filter(Project.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    history = json.loads(project.history)

    if data.section not in history:
        history[data.section] = []

    history[data.section].append({"feedback": data.feedback})

    project.history = json.dumps(history)
    db.commit()

    return {"message": "Feedback saved"}


# -------------------------
# COMMENT MODEL
# -------------------------
class CommentModel(BaseModel):
    section: str
    comment: str


@router.post("/{project_id}/comment")
def add_comment(project_id: int, data: CommentModel, token: str, db: Session = Depends(get_db)):
    user = get_current_user(token, db)

    project = db.query(Project).filter(Project.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    history = json.loads(project.history)

    if data.section not in history:
        history[data.section] = []

    history[data.section].append({"comment": data.comment})

    project.history = json.dumps(history)
    db.commit()

    return {"message": "Comment saved"}


# -------------------------
# GET CONTENT
# -------------------------
@router.get("/{project_id}/get-content")
def get_content(project_id: int, token: str, db: Session = Depends(get_db)):
    user = get_current_user(token, db)

    project = db.query(Project).filter(Project.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    return json.loads(project.content)


# -------------------------
# EXPORT DOCX
# -------------------------
@router.get("/{project_id}/export-docx")
def export_docx(project_id: int, token: str, db: Session = Depends(get_db)):
    user = get_current_user(token, db)

    project = db.query(Project).filter(Project.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    content = json.loads(project.content)

    doc = Document()
    doc.add_heading(project.title, level=1)

    for section, text in content.items():
        doc.add_heading(section, level=2)
        doc.add_paragraph(text)

    os.makedirs("exports", exist_ok=True)
    filepath = f"exports/project_{project_id}.docx"
    doc.save(filepath)

    return FileResponse(filepath, filename=f"project_{project_id}.docx")


# -------------------------
# EXPORT PPTX
# -------------------------
@router.get("/{project_id}/export-pptx")
def export_pptx(project_id: int, token: str, db: Session = Depends(get_db)):
    user = get_current_user(token, db)

    project = db.query(Project).filter(Project.id == project_id).first()
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")

    content = json.loads(project.content)

    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = project.title

    for section, text in content.items():
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = section
        slide.placeholders[1].text = text[:2000]

    os.makedirs("exports", exist_ok=True)
    filepath = f"exports/project_{project_id}.pptx"
    prs.save(filepath)

    return FileResponse(filepath, filename=f"project_{project_id}.pptx")
